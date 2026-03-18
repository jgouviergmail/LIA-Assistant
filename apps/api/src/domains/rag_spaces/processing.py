"""
RAG document processing pipeline.

Background task that extracts text from uploaded documents, splits into
chunks, generates embeddings via TrackedOpenAIEmbeddings, and persists
vector-indexed chunks to the rag_chunks table.

Launched via safe_fire_and_forget after document upload.

Phase: evolution — RAG Spaces (User Knowledge Documents)
Created: 2026-03-14
"""

from __future__ import annotations

import asyncio
import time
from pathlib import Path
from uuid import UUID

import tiktoken
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sqlalchemy.ext.asyncio import AsyncSession

from src.core.config import settings
from src.domains.rag_spaces.embedding import get_rag_embeddings
from src.domains.rag_spaces.models import RAGChunk, RAGDocument, RAGDocumentStatus
from src.domains.rag_spaces.repository import RAGChunkRepository, RAGDocumentRepository
from src.infrastructure.database.session import get_db_context
from src.infrastructure.llm.embedding_context import (
    clear_embedding_context,
    set_embedding_context,
)
from src.infrastructure.observability.logging import get_logger
from src.infrastructure.observability.metrics_rag_spaces import (
    rag_document_chunks_total,
    rag_document_processing_duration_seconds,
    rag_document_upload_size_bytes,
    rag_documents_processed_total,
    rag_documents_total_count,
    rag_embedding_tokens_total,
)

logger = get_logger(__name__)

# Maximum number of chunks to embed in a single API call.
# OpenAI Embeddings API has ~4MB request body limit; 100 chunks of ~1000 chars
# is well within that limit while still being efficient.
EMBEDDING_BATCH_SIZE = 100


# ============================================================================
# Text Extraction
# ============================================================================


def extract_text_plain(file_path: Path) -> str:
    """Extract text from plain text or markdown files."""
    return file_path.read_text(encoding="utf-8", errors="replace")


def extract_text_pdf(file_path: Path) -> str:
    """Extract text from PDF using PyMuPDF (fitz)."""
    import fitz  # type: ignore[import-untyped]  # PyMuPDF

    text_parts: list[str] = []
    with fitz.open(str(file_path)) as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


def extract_text_docx(file_path: Path) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document

    doc = Document(str(file_path))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def extract_text_pptx(file_path: Path) -> str:
    """Extract text from PPTX slides, tables, and speaker notes."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        text_parts.append(text)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        text_parts.append("\t".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                text_parts.append(notes)
    return "\n".join(text_parts)


def extract_text_xlsx(file_path: Path) -> str:
    """Extract text from XLSX spreadsheet cells, sheet by sheet."""
    from openpyxl import load_workbook  # type: ignore[import-untyped]

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    text_parts: list[str] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
        if rows:
            text_parts.append(f"[{sheet.title}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(text_parts)


def extract_text_csv(file_path: Path) -> str:
    """Extract text from CSV file as tab-separated rows."""
    import csv

    text_parts: list[str] = []
    with file_path.open(encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        for row in reader:
            if any(cell.strip() for cell in row):
                text_parts.append("\t".join(row))
    return "\n".join(text_parts)


def extract_text_rtf(file_path: Path) -> str:
    """Extract text from RTF using striprtf."""
    from striprtf.striprtf import rtf_to_text  # type: ignore[import-untyped]

    raw = file_path.read_text(encoding="utf-8", errors="replace")
    return str(rtf_to_text(raw))


def extract_text_html(file_path: Path) -> str:
    """Extract readable text from HTML using markdownify."""
    from markdownify import markdownify

    raw_html = file_path.read_text(encoding="utf-8", errors="replace")
    return markdownify(raw_html, strip=["img", "script", "style"])


def _odf_extract_text(node: object) -> str:
    """Recursively extract text from an ODF node tree.

    Handles nested ``<text:span>`` elements within ``<text:p>`` paragraphs.
    Text nodes (nodeType == 3) carry data; element nodes (nodeType == 1)
    contain children that must be traversed.
    """
    if getattr(node, "nodeType", None) == 3:  # Text node
        return getattr(node, "data", "")
    parts: list[str] = []
    for child in getattr(node, "childNodes", []):
        parts.append(_odf_extract_text(child))
    return "".join(parts)


def extract_text_odt(file_path: Path) -> str:
    """Extract text from ODT (OpenDocument Text)."""
    from odf import text as odf_text  # type: ignore[import-untyped]
    from odf.opendocument import load  # type: ignore[import-untyped]

    doc = load(str(file_path))
    text_parts: list[str] = []
    for paragraph in doc.getElementsByType(odf_text.P):
        content = _odf_extract_text(paragraph)
        if content.strip():
            text_parts.append(content)
    return "\n".join(text_parts)


def extract_text_ods(file_path: Path) -> str:
    """Extract text from ODS (OpenDocument Spreadsheet)."""
    from odf.opendocument import load
    from odf.table import Table, TableCell, TableRow  # type: ignore[import-untyped]
    from odf.text import P  # type: ignore[import-untyped]

    doc = load(str(file_path))
    text_parts: list[str] = []
    for table in doc.getElementsByType(Table):
        table_name = table.getAttribute("name") or "Sheet"
        rows: list[str] = []
        for row in table.getElementsByType(TableRow):
            cells: list[str] = []
            for cell in row.getElementsByType(TableCell):
                cell_texts = []
                for p in cell.getElementsByType(P):
                    cell_texts.append(_odf_extract_text(p))
                cells.append(" ".join(cell_texts))
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
        if rows:
            text_parts.append(f"[{table_name}]\n" + "\n".join(rows))
    return "\n\n".join(text_parts)


def extract_text_odp(file_path: Path) -> str:
    """Extract text from ODP (OpenDocument Presentation)."""
    from odf.draw import Frame, Page  # type: ignore[import-untyped]
    from odf.opendocument import load
    from odf.text import P

    doc = load(str(file_path))
    text_parts: list[str] = []
    for page in doc.getElementsByType(Page):
        for frame in page.getElementsByType(Frame):
            for p in frame.getElementsByType(P):
                content = _odf_extract_text(p)
                if content.strip():
                    text_parts.append(content)
    return "\n".join(text_parts)


def extract_text_epub(file_path: Path) -> str:
    """Extract chapter text from EPUB in spine (reading) order."""
    import warnings

    import ebooklib  # type: ignore[import-untyped]
    from ebooklib import epub
    from markdownify import markdownify

    with warnings.catch_warnings():
        warnings.simplefilter("ignore")
        book = epub.read_epub(str(file_path))

    text_parts: list[str] = []
    for item_id, _linear in book.spine:
        item = book.get_item_with_id(item_id)
        if item is None:
            continue
        if item.get_type() != ebooklib.ITEM_DOCUMENT:
            continue
        html_content = item.get_content().decode("utf-8", errors="replace")
        chapter_text = markdownify(html_content, strip=["img", "script", "style"])
        if chapter_text.strip():
            text_parts.append(chapter_text.strip())
    return "\n\n".join(text_parts)


def extract_text_json(file_path: Path) -> str:
    """Extract JSON content as pretty-printed text."""
    import json

    raw = file_path.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(raw)
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return raw


def extract_text_xml(file_path: Path) -> str:
    """Extract XML content as formatted text using defusedxml (XXE-safe)."""
    import xml.etree.ElementTree as ET

    import defusedxml.ElementTree as DefusedET  # type: ignore[import-untyped]

    try:
        tree = DefusedET.parse(str(file_path))
        root = tree.getroot()
        ET.indent(root)
        return ET.tostring(root, encoding="unicode")
    except ET.ParseError:
        return file_path.read_text(encoding="utf-8", errors="replace")


def extract_text(file_path: Path, content_type: str) -> str:
    """
    Extract text from a file based on its MIME type.

    Args:
        file_path: Path to the file on disk
        content_type: MIME type of the file

    Returns:
        Extracted text content

    Raises:
        ValueError: If content type is not supported
    """
    extractors = {
        "text/plain": extract_text_plain,
        "text/markdown": extract_text_plain,
        "application/pdf": extract_text_pdf,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": (
            extract_text_docx
        ),
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": (
            extract_text_pptx
        ),
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": extract_text_xlsx,
        "text/csv": extract_text_csv,
        "application/rtf": extract_text_rtf,
        "text/html": extract_text_html,
        "application/vnd.oasis.opendocument.text": extract_text_odt,
        "application/vnd.oasis.opendocument.spreadsheet": extract_text_ods,
        "application/vnd.oasis.opendocument.presentation": extract_text_odp,
        "application/epub+zip": extract_text_epub,
        "application/json": extract_text_json,
        "application/xml": extract_text_xml,
        "text/xml": extract_text_xml,
    }
    extractor = extractors.get(content_type)
    if not extractor:
        msg = f"Unsupported content type: {content_type}"
        raise ValueError(msg)

    return extractor(file_path)


# ============================================================================
# Processing Pipeline
# ============================================================================


async def process_document(
    document_id: UUID,
    space_id: UUID,
    user_id: UUID,
    filename: str,
    original_filename: str,
    content_type: str,
) -> None:
    """
    Background task to process an uploaded document.

    Pipeline:
    1. Extract text from file
    2. Split into chunks
    3. Generate embeddings (batch)
    4. Persist chunks to rag_chunks table
    5. Update document status to ready

    This function creates its own DB session via get_db_context()
    and sets the embedding tracking context for cost attribution.
    """
    start_time = time.time()
    logger.info(
        "rag_document_processing_started",
        document_id=str(document_id),
        space_id=str(space_id),
        original_filename=original_filename,
        content_type=content_type,
    )

    # Set embedding tracking context for cost attribution
    set_embedding_context(
        user_id=str(user_id),
        session_id=f"rag_index_{document_id}",
        conversation_id=None,
    )

    # Track previous status for gauge transitions (processing or reindexing)
    previous_status = RAGDocumentStatus.PROCESSING

    try:
        async with get_db_context() as db:
            doc_repo = RAGDocumentRepository(db)
            chunk_repo = RAGChunkRepository(db)

            # Fetch document
            document = await doc_repo.get_by_id(document_id)
            if not document:
                logger.error(
                    "rag_document_not_found_for_processing",
                    document_id=str(document_id),
                )
                return

            previous_status = document.status

            # 1. Extract text
            file_path = (
                Path(settings.rag_spaces_storage_path) / str(user_id) / str(space_id) / filename
            )
            if not file_path.exists():
                await _mark_document_error(doc_repo, document, db, "File not found on disk")
                return

            try:
                text = await asyncio.to_thread(extract_text, file_path, content_type)
            except Exception as e:
                await _mark_document_error(doc_repo, document, db, f"Text extraction failed: {e}")
                return

            if not text.strip():
                await _mark_document_error(doc_repo, document, db, "No text content extracted")
                return

            # 2. Split into chunks
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=settings.rag_spaces_chunk_size,
                chunk_overlap=settings.rag_spaces_chunk_overlap,
                length_function=len,
                separators=["\n\n", "\n", ". ", " ", ""],
            )
            chunk_texts = splitter.split_text(text)

            if not chunk_texts:
                await _mark_document_error(
                    doc_repo, document, db, "No chunks produced after splitting"
                )
                return

            max_chunks = settings.rag_spaces_max_chunks_per_document
            if len(chunk_texts) > max_chunks:
                await _mark_document_error(
                    doc_repo,
                    document,
                    db,
                    f"Document produces {len(chunk_texts)} chunks, exceeding limit of {max_chunks}",
                )
                return

            logger.debug(
                "rag_document_chunks_created",
                document_id=str(document_id),
                chunk_count=len(chunk_texts),
            )

            # 3. Count tokens before embedding (for cost tracking on document)
            try:
                encoding = tiktoken.get_encoding("cl100k_base")
                total_embedding_tokens = sum(len(encoding.encode(t)) for t in chunk_texts if t)
            except Exception:
                total_embedding_tokens = sum(len(t) // 4 for t in chunk_texts if t)

            # 4. Generate embeddings in batches to avoid API size limits
            embeddings_model = get_rag_embeddings()
            embedding_vectors: list[list[float]] = []
            for batch_start in range(0, len(chunk_texts), EMBEDDING_BATCH_SIZE):
                batch = chunk_texts[batch_start : batch_start + EMBEDDING_BATCH_SIZE]
                batch_vectors = await embeddings_model.aembed_documents(batch)
                embedding_vectors.extend(batch_vectors)

            # 5. Build chunk objects
            model_name = settings.rag_spaces_embedding_model
            chunk_objects = []
            for idx, (chunk_text, embedding_vec) in enumerate(
                zip(chunk_texts, embedding_vectors, strict=True)
            ):
                chunk_objects.append(
                    RAGChunk(
                        document_id=document_id,
                        space_id=space_id,
                        user_id=user_id,
                        chunk_index=idx,
                        content=chunk_text,
                        embedding=embedding_vec,
                        embedding_model=model_name,
                        metadata_={
                            "original_filename": original_filename,
                            "content_type": content_type,
                            "chunk_index": idx,
                            "total_chunks": len(chunk_texts),
                        },
                    )
                )

            # 6. Bulk insert chunks
            await chunk_repo.bulk_create_chunks(chunk_objects)

            # 7. Calculate embedding cost (reuse shared pricing logic)
            from src.infrastructure.cache.pricing_cache import get_cached_usd_eur_rate
            from src.infrastructure.llm.tracked_embeddings import estimate_embedding_cost_sync

            embedding_cost_usd = estimate_embedding_cost_sync(model_name, total_embedding_tokens)
            embedding_cost_eur = round(embedding_cost_usd * get_cached_usd_eur_rate(), 6)

            # 8. Update document status with token/cost tracking
            await doc_repo.update(
                document,
                {
                    "status": RAGDocumentStatus.READY,
                    "chunk_count": len(chunk_objects),
                    "embedding_model": model_name,
                    "embedding_tokens": total_embedding_tokens,
                    "embedding_cost_eur": embedding_cost_eur,
                    "error_message": None,
                },
            )
            await db.commit()

            duration = time.time() - start_time

            # Prometheus metrics
            rag_documents_processed_total.labels(status="success").inc()
            rag_documents_total_count.labels(status=previous_status).dec()
            rag_documents_total_count.labels(status="ready").inc()
            rag_document_processing_duration_seconds.observe(duration)
            rag_document_chunks_total.observe(len(chunk_objects))
            rag_document_upload_size_bytes.labels(content_type=content_type).observe(
                document.file_size
            )
            rag_embedding_tokens_total.labels(operation="index").inc(total_embedding_tokens)

            logger.info(
                "rag_document_processing_complete",
                document_id=str(document_id),
                chunk_count=len(chunk_objects),
                duration_seconds=round(duration, 2),
                embedding_model=model_name,
            )

    except Exception as e:
        logger.error(
            "rag_document_processing_failed",
            document_id=str(document_id),
            error=str(e),
            exc_info=True,
        )
        # Try to mark the document as error (also updates gauge metrics)
        try:
            async with get_db_context() as db:
                doc_repo = RAGDocumentRepository(db)
                document = await doc_repo.get_by_id(document_id)
                if document:
                    await _mark_document_error(doc_repo, document, db, f"Processing failed: {e}")
        except Exception:
            logger.error(
                "rag_document_error_status_update_failed",
                document_id=str(document_id),
            )

    finally:
        clear_embedding_context()


async def _mark_document_error(
    doc_repo: RAGDocumentRepository,
    document: RAGDocument,
    db: AsyncSession,
    error_message: str,
) -> None:
    """Mark a document as error with a message and update gauge metrics."""
    previous_status = document.status
    await doc_repo.update(
        document,
        {
            "status": RAGDocumentStatus.ERROR,
            "error_message": error_message,
        },
    )
    await db.commit()
    rag_documents_processed_total.labels(status="error").inc()
    rag_documents_total_count.labels(status=previous_status).dec()
    rag_documents_total_count.labels(status="error").inc()

    logger.warning(
        "rag_document_processing_error",
        document_id=str(document.id),
        error_message=error_message,
    )
