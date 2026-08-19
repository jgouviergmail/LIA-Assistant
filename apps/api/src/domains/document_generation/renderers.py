"""Pure renderers: structured content -> document bytes (ADR-226).

Every renderer is a pure function (content in, bytes out) so it is unit-tested
without I/O; CPU-bound rendering is offloaded with ``asyncio.to_thread`` by the
CALLER (service layer). Heavy libraries (openpyxl, python-docx, python-pptx,
PyMuPDF) are imported lazily inside their renderer so importing this module
stays cheap. The registry is completeness-asserted at import (ADR-085).
"""

from __future__ import annotations

import csv
import io
import re
from collections.abc import Callable

from src.domains.document_generation.sanitize import neutralize_formula
from src.domains.document_generation.schemas import (
    DocumentContent,
    DocumentType,
    SectionBlock,
    SectionedContent,
    SlideContent,
    TableSheet,
    TabularContent,
)

DOCUMENT_MIME_TYPES: dict[DocumentType, str] = {
    DocumentType.CSV: "text/csv",
    DocumentType.XLSX: ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
    DocumentType.DOCX: ("application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    DocumentType.PPTX: (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ),
    DocumentType.PDF: "application/pdf",
    DocumentType.MD: "text/markdown",
    DocumentType.TXT: "text/plain",
}

DOCUMENT_EXTENSIONS: dict[DocumentType, str] = {
    DocumentType.CSV: "csv",
    DocumentType.XLSX: "xlsx",
    DocumentType.DOCX: "docx",
    DocumentType.PPTX: "pptx",
    DocumentType.PDF: "pdf",
    DocumentType.MD: "md",
    DocumentType.TXT: "txt",
}


# ---------------------------------------------------------------------------
# Text family: csv / md / txt
# ---------------------------------------------------------------------------


def _render_csv(content: DocumentContent) -> bytes:
    if not isinstance(content, TabularContent):
        raise ValueError("csv rendering requires TabularContent")
    sheet = content.sheets[0]
    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow([neutralize_formula(h) for h in sheet.headers])
    for row in sheet.rows:
        writer.writerow([neutralize_formula(cell) for cell in row])
    # utf-8-sig: Excel needs the BOM to detect UTF-8 (probe 2026-08-17).
    return buf.getvalue().encode("utf-8-sig")


def _md_table(table: TableSheet) -> list[str]:
    header = "| " + " | ".join(table.headers) + " |"
    rule = "| " + " | ".join("---" for _ in table.headers) + " |"
    rows = ["| " + " | ".join(row) + " |" for row in table.rows]
    return [header, rule, *rows]


def _md_block(block: SectionBlock) -> list[str]:
    if block.kind == "heading":
        # The document title owns "#"; content headings start at "##" even
        # when the LLM says level 1 — same shift the PDF renderer applies.
        return [f"{'#' * max(block.level, 2)} {block.text}"]
    if block.kind == "paragraph":
        return [block.text]
    if block.kind == "bullets":
        return [f"- {item}" for item in block.items]
    if block.table is not None:
        return _md_table(block.table)
    return []


def _render_md(content: DocumentContent) -> bytes:
    if not isinstance(content, SectionedContent):
        raise ValueError("md rendering requires SectionedContent")
    lines: list[str] = [f"# {content.title}", ""]
    for block in content.blocks:
        lines.extend(_md_block(block))
        lines.append("")
    return "\n".join(lines).encode("utf-8")


def _render_txt(content: DocumentContent) -> bytes:
    if not isinstance(content, SectionedContent):
        raise ValueError("txt rendering requires SectionedContent")
    lines: list[str] = [content.title, "=" * len(content.title), ""]
    for block in content.blocks:
        if block.kind == "heading":
            lines.extend([block.text, "-" * len(block.text)])
        elif block.kind == "paragraph":
            lines.append(block.text)
        elif block.kind == "bullets":
            lines.extend(f"  * {item}" for item in block.items)
        elif block.table is not None:
            lines.append(" / ".join(block.table.headers))
            lines.extend(" / ".join(row) for row in block.table.rows)
        lines.append("")
    return "\n".join(lines).encode("utf-8")


# ---------------------------------------------------------------------------
# Office family: xlsx / docx / pptx
# ---------------------------------------------------------------------------

_XLSX_TITLE_FORBIDDEN = re.compile(r"[\[\]:*?/\\]")  # openpyxl rejects these
_XLSX_TITLE_MAX = 31  # Excel's hard sheet-title limit


def _xlsx_sheet_title(name: str, index: int, used: set[str]) -> str:
    """Sanitize an LLM-suggested worksheet title for openpyxl.

    Strips the characters openpyxl rejects, enforces Excel's 31-char limit,
    falls back to ``Sheet{n}`` when empty, deduplicates with a numeric suffix.

    Args:
        name: Suggested sheet name.
        index: Zero-based sheet position (for the fallback name).
        used: Titles already taken in this workbook (mutated in place).

    Returns:
        A unique, openpyxl-legal worksheet title.
    """
    cleaned = _XLSX_TITLE_FORBIDDEN.sub("_", name).strip()[:_XLSX_TITLE_MAX]
    cleaned = cleaned or f"Sheet{index + 1}"
    candidate = cleaned
    suffix = 2
    while candidate in used:
        tail = f" {suffix}"
        candidate = f"{cleaned[: _XLSX_TITLE_MAX - len(tail)]}{tail}"
        suffix += 1
    used.add(candidate)
    return candidate


def _render_xlsx(content: DocumentContent) -> bytes:
    if not isinstance(content, TabularContent):
        raise ValueError("xlsx rendering requires TabularContent")
    import openpyxl
    from openpyxl.styles import Font
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    default_sheet = wb.active
    used_titles: set[str] = set()
    for index, sheet in enumerate(content.sheets):
        ws = default_sheet if index == 0 else wb.create_sheet()
        ws.title = _xlsx_sheet_title(sheet.name, index, used_titles)
        ws.append([neutralize_formula(header) for header in sheet.headers])
        for cell in ws[1]:
            cell.font = Font(bold=True)
        for row in sheet.rows:
            ws.append([neutralize_formula(value) for value in row])
        for col_index, header in enumerate(sheet.headers, start=1):
            widths = [len(header)] + [
                len(row[col_index - 1]) for row in sheet.rows if len(row) >= col_index
            ]
            ws.column_dimensions[get_column_letter(col_index)].width = min(max(widths) + 2, 60)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _render_docx(content: DocumentContent) -> bytes:
    if not isinstance(content, SectionedContent):
        raise ValueError("docx rendering requires SectionedContent")
    import docx

    document = docx.Document()
    document.add_heading(content.title, level=1)
    for block in content.blocks:
        if block.kind == "heading":
            document.add_heading(block.text, level=min(block.level, 4))
        elif block.kind == "paragraph":
            document.add_paragraph(block.text)
        elif block.kind == "bullets":
            for item in block.items:
                document.add_paragraph(item, style="List Bullet")
        elif block.table is not None:
            table = document.add_table(
                rows=len(block.table.rows) + 1, cols=len(block.table.headers)
            )
            table.style = "Light Grid Accent 1"
            for col, header in enumerate(block.table.headers):
                table.cell(0, col).text = header
            for row_index, row in enumerate(block.table.rows, start=1):
                for col, value in enumerate(row[: len(block.table.headers)]):
                    table.cell(row_index, col).text = value
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _render_pptx(content: DocumentContent) -> bytes:
    if not isinstance(content, SlideContent):
        raise ValueError("pptx rendering requires SlideContent")
    import pptx

    presentation = pptx.Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = content.title
    for slide_spec in content.slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_spec.title
        body = slide.placeholders[1].text_frame
        for index, bullet in enumerate(slide_spec.bullets):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet
        if slide_spec.notes:
            slide.notes_slide.notes_text_frame.text = slide_spec.notes
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PDF (PyMuPDF Story: escaped HTML -> paged A4)
# ---------------------------------------------------------------------------


def _pdf_html(content: SectionedContent) -> str:
    """Build the escaped HTML the Story engine lays out."""
    from html import escape

    parts: list[str] = [f"<h1>{escape(content.title)}</h1>"]
    for block in content.blocks:
        if block.kind == "heading":
            # The title owns h1; content headings shift one level down,
            # mirroring the markdown renderer.
            level = min(block.level + 1, 5)
            parts.append(f"<h{level}>{escape(block.text)}</h{level}>")
        elif block.kind == "paragraph":
            parts.append(f"<p>{escape(block.text)}</p>")
        elif block.kind == "bullets":
            items = "".join(f"<li>{escape(item)}</li>" for item in block.items)
            parts.append(f"<ul>{items}</ul>")
        elif block.table is not None:
            head = "".join(f"<th>{escape(header)}</th>" for header in block.table.headers)
            body = "".join(
                "<tr>" + "".join(f"<td>{escape(value)}</td>" for value in row) + "</tr>"
                for row in block.table.rows
            )
            parts.append(f"<table><tr>{head}</tr>{body}</table>")
    return "".join(parts)


def _render_pdf(content: DocumentContent) -> bytes:
    if not isinstance(content, SectionedContent):
        raise ValueError("pdf rendering requires SectionedContent")
    import fitz  # type: ignore[import-untyped]  # PyMuPDF

    story = fitz.Story(html=_pdf_html(content))
    buf = io.BytesIO()
    writer = fitz.DocumentWriter(buf)
    mediabox = fitz.paper_rect("a4")
    where = mediabox + (36, 36, -36, -36)
    more = True
    while more:
        device = writer.begin_page(mediabox)
        more, _ = story.place(where)
        story.draw(device)
        writer.end_page()
    writer.close()
    return buf.getvalue()


RENDERERS: dict[DocumentType, Callable[[DocumentContent], bytes]] = {
    DocumentType.CSV: _render_csv,
    DocumentType.MD: _render_md,
    DocumentType.TXT: _render_txt,
    DocumentType.XLSX: _render_xlsx,
    DocumentType.DOCX: _render_docx,
    DocumentType.PPTX: _render_pptx,
    DocumentType.PDF: _render_pdf,
}

# Boot-time completeness (ADR-085): a partial renderer map refuses to import.
assert set(RENDERERS) == set(DocumentType), "RENDERERS must cover every DocumentType"


def render_document(doc_type: DocumentType, content: DocumentContent) -> bytes:
    """Render structured content into final document bytes.

    Args:
        doc_type: Target format.
        content: Validated content matching ``SCHEMA_BY_DOC_TYPE[doc_type]``.

    Returns:
        The rendered file bytes.

    Raises:
        ValueError: When the content model does not match the format family.
    """
    return RENDERERS[doc_type](content)
