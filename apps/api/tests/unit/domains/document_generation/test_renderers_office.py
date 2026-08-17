"""Office renderers round-trip through their own readers (ADR-226).

The oracles are the same libraries the RAG extractors already embed
(openpyxl / python-docx / python-pptx), so "what we write is what a reader
sees" is asserted, not assumed.
"""

import io

import docx
import openpyxl
import pptx
import pytest

from src.domains.document_generation.renderers import render_document
from src.domains.document_generation.schemas import (
    DocumentType,
    SectionBlock,
    SectionedContent,
    Slide,
    SlideContent,
    TableSheet,
    TabularContent,
)


@pytest.mark.unit
class TestXlsxRenderer:
    """xlsx: round-trip, neutralization, sheet-title safety."""

    def test_round_trip_and_formula_neutralized(self) -> None:
        content = TabularContent(
            filename_stem="data",
            title="Data",
            sheets=[
                TableSheet(name="Feuille 1", headers=["a", "b"], rows=[["1", "=2+2"]]),
                TableSheet(name="Feuille 2", headers=["c"], rows=[["x"]]),
            ],
        )
        data = render_document(DocumentType.XLSX, content)
        wb = openpyxl.load_workbook(io.BytesIO(data))
        assert wb.sheetnames == ["Feuille 1", "Feuille 2"]
        ws = wb["Feuille 1"]
        assert ws["A1"].value == "a"
        assert ws["B2"].value == "'=2+2"
        assert ws["B2"].data_type != "f"  # the probe-proven injection stays closed

    def test_negative_numbers_survive_untouched(self) -> None:
        content = TabularContent(
            filename_stem="deltas",
            title="Deltas",
            sheets=[TableSheet(name="D", headers=["delta"], rows=[["-5.2"]])],
        )
        wb = openpyxl.load_workbook(io.BytesIO(render_document(DocumentType.XLSX, content)))
        assert wb.active["A2"].value == "-5.2"

    def test_sheet_titles_sanitized_and_deduplicated(self) -> None:
        # openpyxl REJECTS []:*?/\ in titles; the LLM can produce both invalid
        # characters and duplicate names — the renderer must survive both.
        content = TabularContent(
            filename_stem="data",
            title="Data",
            sheets=[
                TableSheet(name="Q1/Q2 [draft]", headers=["a"], rows=[["1"]]),
                TableSheet(name="Q1/Q2 [draft]", headers=["b"], rows=[["2"]]),
                TableSheet(name="", headers=["c"], rows=[["3"]]),
            ],
        )
        wb = openpyxl.load_workbook(io.BytesIO(render_document(DocumentType.XLSX, content)))
        assert len(wb.sheetnames) == 3
        assert len(set(wb.sheetnames)) == 3  # deduplicated
        for title in wb.sheetnames:
            assert not set(title) & set("[]:*?/\\")  # sanitized

    def test_requires_tabular_content(self) -> None:
        content = SectionedContent(
            filename_stem="x",
            title="T",
            blocks=[SectionBlock(kind="paragraph", text="p")],
        )
        with pytest.raises(ValueError, match="TabularContent"):
            render_document(DocumentType.XLSX, content)


@pytest.mark.unit
class TestDocxRenderer:
    """docx: every block kind lands where python-docx can read it back."""

    def test_round_trip_blocks(self) -> None:
        content = SectionedContent(
            filename_stem="rapport",
            title="Rapport",
            blocks=[
                SectionBlock(kind="heading", level=2, text="Partie 1"),
                SectionBlock(kind="paragraph", text="Texte accentué éàü."),
                SectionBlock(kind="bullets", items=["un", "deux"]),
                SectionBlock(
                    kind="table",
                    table=TableSheet(name="T", headers=["k"], rows=[["v"]]),
                ),
            ],
        )
        data = render_document(DocumentType.DOCX, content)
        document = docx.Document(io.BytesIO(data))
        texts = [p.text for p in document.paragraphs]
        assert "Rapport" in texts
        assert "Partie 1" in texts
        assert "Texte accentué éàü." in texts
        assert "un" in texts and "deux" in texts
        assert document.tables
        assert document.tables[0].cell(0, 0).text == "k"
        assert document.tables[0].cell(1, 0).text == "v"


@pytest.mark.unit
class TestPptxRenderer:
    """pptx: title slide + one slide per spec, bullets and notes intact."""

    def test_round_trip_slides_and_notes(self) -> None:
        content = SlideContent(
            filename_stem="alsace",
            title="L'Alsace",
            slides=[
                Slide(
                    title="Géographie",
                    bullets=["Rhin", "Vosges"],
                    notes="parler lentement",
                )
            ],
        )
        data = render_document(DocumentType.PPTX, content)
        presentation = pptx.Presentation(io.BytesIO(data))
        assert len(presentation.slides) == 2  # title slide + 1 content slide
        all_text = [
            paragraph.text
            for slide in presentation.slides
            for shape in slide.shapes
            if shape.has_text_frame
            for paragraph in shape.text_frame.paragraphs
        ]
        assert "L'Alsace" in all_text
        assert "Géographie" in all_text
        assert "Rhin" in all_text and "Vosges" in all_text
        notes_slide = presentation.slides[1].notes_slide
        assert "parler lentement" in notes_slide.notes_text_frame.text

    def test_slide_without_bullets_is_valid(self) -> None:
        content = SlideContent(
            filename_stem="x",
            title="T",
            slides=[Slide(title="Vide")],
        )
        data = render_document(DocumentType.PPTX, content)
        presentation = pptx.Presentation(io.BytesIO(data))
        assert len(presentation.slides) == 2
